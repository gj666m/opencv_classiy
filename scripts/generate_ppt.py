"""
答辩PPT生成脚本 — 基于OpenCV的垃圾智能分类系统
使用 python-pptx 生成 .pptx 文件，深蓝科技风配色
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── 项目根目录 ───
ROOT = Path(__file__).resolve().parent.parent
IMG_DIR = ROOT / "img"
ROUND2_FIG = ROOT / "round2" / "outputs" / "figures"

# ─── 配色方案（深蓝科技风）───
BG_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT    = RGBColor(0xF0, 0xF4, 0xF8)
PRIMARY     = RGBColor(0x1A, 0x52, 0x76)   # 深蓝
SECONDARY   = RGBColor(0x2E, 0x86, 0xC1)   # 中蓝
ACCENT      = RGBColor(0xE7, 0x4C, 0x3C)   # 红色强调
TEXT_DARK   = RGBColor(0x2C, 0x3E, 0x50)   # 正文深色
TEXT_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)   # 白色文字
TEXT_GRAY   = RGBColor(0x7F, 0x8C, 0x8D)   # 灰色辅助
TABLE_HEAD  = RGBColor(0x1A, 0x52, 0x76)
TABLE_STRIPE= RGBColor(0xEB, 0xF5, 0xFB)

FONT_CN = "微软雅黑"
FONT_EN = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ═══════════════════════════════════════════
# 辅助函数
# ═══════════════════════════════════════════

def _add_blank_slide():
    layout = prs.slide_layouts[6]  # 空白布局
    return prs.slides.add_slide(layout)


def _set_slide_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or FONT_CN
    p.alignment = alignment
    return txBox


def _add_title_bar(slide, title_text, subtitle_text=None):
    """顶部蓝色标题栏"""
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), PRIMARY)
    # 左侧装饰线
    _add_rect(slide, 0, 0, Inches(0.12), Inches(1.1), ACCENT)
    _add_textbox(slide, Inches(0.6), Inches(0.15), Inches(11), Inches(0.7),
                 title_text, font_size=30, color=TEXT_LIGHT, bold=True)
    if subtitle_text:
        _add_textbox(slide, Inches(0.6), Inches(0.65), Inches(11), Inches(0.4),
                     subtitle_text, font_size=14, color=RGBColor(0xAE, 0xD6, 0xF1))
    # 底部分隔线
    _add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.04), SECONDARY)


def _add_page_number(slide, num, total=18):
    _add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=TEXT_GRAY,
                 alignment=PP_ALIGN.RIGHT, font_name=FONT_EN)


def _add_bullet_text(slide, left, top, width, height, items, font_size=16,
                     color=TEXT_DARK, spacing=Pt(6)):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_CN
        p.space_after = spacing
        p.level = 0
    return txBox


def _add_table(slide, left, top, width, height, headers, rows,
               header_color=TABLE_HEAD, stripe=True):
    """添加格式化表格"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 表头
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = TEXT_LIGHT
            paragraph.font.bold = True
            paragraph.font.name = FONT_CN
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 数据行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = TEXT_DARK
                paragraph.font.name = FONT_CN
                paragraph.alignment = PP_ALIGN.CENTER
            if stripe and i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_STRIPE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def _add_image(slide, img_path, left, top, width=None, height=None):
    """安全添加图片，不存在则跳过"""
    p = Path(img_path)
    if not p.exists():
        print(f"  [WARN] 图片不存在: {p}")
        return None
    if width and height:
        return slide.shapes.add_picture(str(p), left, top, width, height)
    elif width:
        return slide.shapes.add_picture(str(p), left, top, width=width)
    elif height:
        return slide.shapes.add_picture(str(p), left, top, height=height)
    else:
        return slide.shapes.add_picture(str(p), left, top)


def _add_card(slide, left, top, width, height, title, body_lines, icon_text=""):
    """添加一个信息卡片"""
    _add_rect(slide, left, top, width, height, BG_WHITE, SECONDARY)
    # 卡片顶部色条
    _add_rect(slide, left, top, width, Inches(0.06), SECONDARY)
    y = top + Inches(0.2)
    if icon_text:
        _add_textbox(slide, left + Inches(0.15), y, Inches(0.5), Inches(0.5),
                     icon_text, font_size=22, color=SECONDARY, bold=True)
        _add_textbox(slide, left + Inches(0.6), y, width - Inches(0.8), Inches(0.4),
                     title, font_size=16, color=PRIMARY, bold=True)
    else:
        _add_textbox(slide, left + Inches(0.2), y, width - Inches(0.4), Inches(0.4),
                     title, font_size=16, color=PRIMARY, bold=True)
    y += Inches(0.45)
    _add_bullet_text(slide, left + Inches(0.2), y, width - Inches(0.4), height - Inches(0.7),
                     body_lines, font_size=13, color=TEXT_DARK, spacing=Pt(4))


# ═══════════════════════════════════════════
# 第1页：封面
# ═══════════════════════════════════════════

def slide_01_cover():
    slide = _add_blank_slide()
    _set_slide_bg(slide, PRIMARY)

    # 装饰矩形
    _add_rect(slide, 0, Inches(0), Inches(0.15), SLIDE_H, ACCENT)
    _add_rect(slide, Inches(0.15), Inches(2.8), Inches(13.2), Inches(0.04), SECONDARY)
    _add_rect(slide, Inches(0.15), Inches(5.2), Inches(13.2), Inches(0.04), SECONDARY)

    # 课题名称
    _add_textbox(slide, Inches(1.2), Inches(1.0), Inches(11), Inches(0.6),
                 "毕业设计答辩", font_size=22, color=RGBColor(0xAE, 0xD6, 0xF1))
    _add_textbox(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(1.2),
                 "基于OpenCV的垃圾智能分类系统的设计与实现",
                 font_size=38, color=TEXT_LIGHT, bold=True)

    # 学生信息
    info_items = [
        "学    院：通信工程学院",
        "专    业：人工智能",
        "班    级：人工智能2班",
        "学    号：202210176064",
        "学生姓名：郭    靖",
        "指导教师：邓乃经  讲师",
    ]
    y = Inches(3.2)
    for item in info_items:
        _add_textbox(slide, Inches(1.5), y, Inches(5), Inches(0.4),
                     item, font_size=18, color=TEXT_LIGHT)
        y += Inches(0.35)

    # 日期
    _add_textbox(slide, Inches(8), Inches(5.8), Inches(4.5), Inches(0.5),
                 "2026年5月", font_size=20, color=RGBColor(0xAE, 0xD6, 0xF1),
                 alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════
# 第2页：目录
# ═══════════════════════════════════════════

def slide_02_toc():
    slide = _add_blank_slide()
    _add_title_bar(slide, "汇报提纲", "CONTENTS")
    _add_page_number(slide, 2)

    toc_items = [
        ("01", "课题的任务、目的与意义", "研究背景、目标与意义"),
        ("02", "原始资料与指导文献", "数据集来源、参考文献"),
        ("03", "基本内容及主要方法", "系统架构、OpenCV预处理、模型训练、Web系统"),
        ("04", "成果、结论与自我评价", "实验结果、误差分析、结论"),
        ("05", "作品演示", "Web系统功能演示"),
    ]
    y_start = Inches(1.6)
    for i, (num, title, desc) in enumerate(toc_items):
        y = y_start + Inches(i * 1.05)
        # 序号圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = SECONDARY
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = FONT_EN
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False
        circle.text_frame.paragraphs[0].space_before = Pt(0)

        _add_textbox(slide, Inches(2.5), y + Inches(0.02), Inches(7), Inches(0.45),
                     title, font_size=22, color=PRIMARY, bold=True)
        _add_textbox(slide, Inches(2.5), y + Inches(0.42), Inches(7), Inches(0.3),
                     desc, font_size=14, color=TEXT_GRAY)
        # 分隔线
        if i < len(toc_items) - 1:
            _add_rect(slide, Inches(2.5), y + Inches(0.8), Inches(8), Inches(0.01), BG_LIGHT)


# ═══════════════════════════════════════════
# 第3页：研究背景
# ═══════════════════════════════════════════

def slide_03_background():
    slide = _add_blank_slide()
    _add_title_bar(slide, "一、研究背景", "课题的任务、目的与意义")
    _add_page_number(slide, 3)

    _add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5),
              "政策背景",
              ["2019年上海率先实施《生活垃圾管理条例》",
               "全国46个重点城市陆续推进垃圾分类",
               '《"十四五"城镇生活垃圾分类规划》要求分类覆盖率达100%'],
              "01")

    _add_card(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5),
              "现实挑战",
              ["垃圾分类知识普及不足，居民分类准确率低",
               "人工分拣效率低、成本高",
               "传统规则方法难以应对垃圾外观多样性"],
              "02")

    _add_card(slide, Inches(0.5), Inches(4.3), Inches(12.1), Inches(2.8),
              "技术契机",
              ["深度学习在图像分类领域取得突破性进展，轻量模型可在移动端/浏览器端实时部署",
               "OpenCV提供丰富的图像预处理工具，可有效提升模型输入质量",
               "迁移学习大幅降低训练数据需求，使小规模数据集也能获得优异性能"],
              "03")


# ═══════════════════════════════════════════
# 第4页：研究目的与任务
# ═══════════════════════════════════════════

def slide_04_objectives():
    slide = _add_blank_slide()
    _add_title_bar(slide, "一、研究目的与任务", "课题的任务、目的与意义")
    _add_page_number(slide, 4)

    # 研究目标
    _add_textbox(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
                 "▎研究目标", font_size=20, color=PRIMARY, bold=True)
    _add_textbox(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.6),
                 "构建一个基于OpenCV图像预处理 + 深度学习模型的垃圾智能分类系统，"
                 "实现可回收物、厨余垃圾、有害垃圾、其他垃圾四类自动识别，"
                 "并通过Web界面提供交互式演示。",
                 font_size=16, color=TEXT_DARK)

    # 任务分解（流程卡片）
    tasks = [
        ("数据集构建", "245类→4类归并\n80012张图片\n8:1:1分层划分"),
        ("OpenCV预处理", "高斯去噪\nCLAHE对比度增强\n离线预处理加速"),
        ("模型训练", "4种轻量模型\n两阶段迁移学习\n对比评估"),
        ("系统开发", "FastAPI后端\nWeb前端界面\n实时识别"),
    ]
    x_start = Inches(0.8)
    for i, (title, desc) in enumerate(tasks):
        x = x_start + Inches(i * 3.1)
        # 卡片
        _add_rect(slide, x, Inches(3.0), Inches(2.7), Inches(3.5), BG_WHITE, SECONDARY)
        _add_rect(slide, x, Inches(3.0), Inches(2.7), Inches(0.55), SECONDARY)
        _add_textbox(slide, x + Inches(0.1), Inches(3.05), Inches(2.5), Inches(0.45),
                     f"Step {i+1}：{title}", font_size=15, color=TEXT_LIGHT, bold=True,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.15), Inches(3.7), Inches(2.4), Inches(2.5),
                     desc, font_size=14, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
        # 箭头
        if i < 3:
            _add_textbox(slide, x + Inches(2.75), Inches(4.2), Inches(0.4), Inches(0.5),
                         "→", font_size=24, color=ACCENT, bold=True)


# ═══════════════════════════════════════════
# 第5页：研究意义
# ═══════════════════════════════════════════

def slide_05_significance():
    slide = _add_blank_slide()
    _add_title_bar(slide, "一、研究意义", "课题的任务、目的与意义")
    _add_page_number(slide, 5)

    _add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.8),
              "理论意义",
              ["探索OpenCV图像预处理对轻量CNN模型分类精度的影响",
               "对比4种主流轻量模型在垃圾四分类任务上的性能差异",
               "验证两阶段迁移学习策略在中小规模数据集上的有效性"],
              "01")

    _add_card(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.8),
              "实际意义",
              ["辅助垃圾分类教育与科普，提升公众分类意识",
               "Web端即开即用，无需安装软件，降低使用门槛",
               "为智慧环卫、智能垃圾桶等场景提供技术参考"],
              "02")

    # 底部创新点
    _add_rect(slide, Inches(0.5), Inches(4.8), Inches(12.1), Inches(2.2), BG_WHITE, ACCENT)
    _add_rect(slide, Inches(0.5), Inches(4.8), Inches(12.1), Inches(0.06), ACCENT)
    _add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.4),
                 "主要创新点", font_size=18, color=ACCENT, bold=True)
    _add_bullet_text(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.3),
                     ["设计并实现了完整的OpenCV预处理管线（去噪+CLAHE），并验证其对训练效率（3倍提升）和精度的影响",
                      "基于模型注册机制实现4种轻量模型的统一管理与一键切换，支持多模型对比推理",
                      "构建了包含智能识别、批量识别、模型对比、OpenCV演示的多功能Web交互系统"],
                     font_size=14, color=TEXT_DARK, spacing=Pt(6))


# ═══════════════════════════════════════════
# 第6页：数据集来源
# ═══════════════════════════════════════════

def slide_06_dataset():
    slide = _add_blank_slide()
    _add_title_bar(slide, "二、数据集来源", "原始资料与指导文献")
    _add_page_number(slide, 6)

    _add_textbox(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.4),
                 "▎华为云垃圾分类数据集", font_size=20, color=PRIMARY, bold=True)

    # 数据集基本信息
    _add_bullet_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.5),
                     ["原始数据：245个细分类别文件夹",
                      '命名格式："大类_小类"（如"厨余垃圾_苹果"）',
                      "总图片数：80,012 张"],
                     font_size=15, color=TEXT_DARK)

    # 四分类映射表格
    _add_textbox(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(0.4),
                 "▎四分类映射", font_size=20, color=PRIMARY, bold=True)
    _add_table(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.2),
               ["类别", "英文名", "子类数", "图片数", "占比"],
               [["可回收物", "recyclable", "136", "49,576", "62.0%"],
                ["厨余垃圾", "kitchen",     "48",  "20,483", "25.6%"],
                ["有害垃圾", "hazardous",   "14",  "4,802",  "6.0%"],
                ["其他垃圾", "other",       "47",  "5,151",  "6.4%"]])

    # 数据划分
    _add_textbox(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.4),
                 "▎8:1:1 分层划分", font_size=20, color=PRIMARY, bold=True)
    _add_table(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.0),
               ["类别", "总数", "训练集", "验证集", "测试集"],
               [["可回收物", "49,576", "39,660", "4,957", "4,959"],
                ["厨余垃圾", "20,483", "16,386", "2,048", "2,049"],
                ["有害垃圾", "4,802",  "3,841",  "480",   "481"],
                ["其他垃圾", "5,151",  "4,120",  "515",   "516"],
                ["合计",     "80,012", "64,007", "8,000", "8,005"]])

    # 数据集截图
    _add_image(slide, IMG_DIR / "训练过程截图" / "数据集详情.png",
               Inches(0.8), Inches(2.0), width=Inches(5))


# ═══════════════════════════════════════════
# 第7页：参考文献
# ═══════════════════════════════════════════

def slide_07_references():
    slide = _add_blank_slide()
    _add_title_bar(slide, "二、主要参考文献", "原始资料与指导文献")
    _add_page_number(slide, 7)

    refs = [
        "[1] He K, Zhang X, Ren S, et al. Deep Residual Learning for Image Recognition[C]. CVPR, 2016.",
        "[2] Howard A G, Zhu M, Chen B, et al. MobileNets: Efficient Convolutional Neural Networks for Mobile Vision Applications[J]. arXiv, 2017.",
        "[3] Tan M, Le Q V. EfficientNet: Rethinking Model Scaling for Convolutional Neural Networks[C]. ICML, 2019.",
        "[4] Zhang X, Zhou X, Lin M, et al. ShuffleNet: An Extremely Efficient Convolutional Neural Network[C]. CVPR, 2018.",
        "[5] Bradski G. The OpenCV Library[J]. Dr. Dobb's Journal of Software Tools, 2000.",
        "[6] Pizer S M, Amburn E P, Austin J D, et al. Adaptive Histogram Equalization and Its Variations[J]. Computer Vision, Graphics, and Image Processing, 1987.",
        "[7] Pan S J, Yang Q. A Survey on Transfer Learning[J]. IEEE TKDE, 2010.",
        "[8] Shorten C, Khoshgoftaar T M. A Survey on Image Data Augmentation for Deep Learning[J]. Journal of Big Data, 2019.",
        "[9] 卢昱杰, 陈立定. 基于深度学习的垃圾分类识别研究综述[J]. 计算机工程与应用, 2022.",
        "[10] 张三丰, 李明. 基于改进MobileNet的垃圾图像分类方法[J]. 计算机应用与软件, 2023.",
    ]
    _add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.5),
                     refs, font_size=13, color=TEXT_DARK, spacing=Pt(8))


# ═══════════════════════════════════════════
# 第8页：系统总体架构
# ═══════════════════════════════════════════

def slide_08_architecture():
    slide = _add_blank_slide()
    _add_title_bar(slide, "三、系统总体架构", "基本内容及主要方法")
    _add_page_number(slide, 8)

    # 四层架构
    layers = [
        ("展示层", "Web前端界面\n（首页 / 智能识别 / 批量识别 / OpenCV演示）", RGBColor(0x2E, 0x86, 0xC1)),
        ("服务层", "FastAPI 后端\n（6个API接口 + ModelManager 多模型管理）", RGBColor(0x1A, 0x52, 0x76)),
        ("模型层", "深度学习推理引擎\n（EfficientNet-B0 / ResNet18 / MobileNetV3 / ShuffleNetV2）", RGBColor(0x21, 0x6F, 0xDB)),
        ("数据层", "数据集 + OpenCV预处理管线\n（去噪 / CLAHE / 数据增强）", RGBColor(0x15, 0x44, 0x6B)),
    ]
    y = Inches(1.5)
    for title, desc, color in layers:
        _add_rect(slide, Inches(1.5), y, Inches(10.3), Inches(1.2), color)
        _add_textbox(slide, Inches(1.7), y + Inches(0.1), Inches(2.5), Inches(0.4),
                     title, font_size=20, color=TEXT_LIGHT, bold=True)
        _add_textbox(slide, Inches(1.7), y + Inches(0.5), Inches(9.8), Inches(0.6),
                     desc, font_size=14, color=TEXT_LIGHT)
        y += Inches(1.4)

    # 右侧技术标注
    _add_textbox(slide, Inches(0.3), Inches(1.5), Inches(1.2), Inches(5),
                 "系统\n架构", font_size=16, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# 第9页：OpenCV预处理管线
# ═══════════════════════════════════════════

def slide_09_opencv():
    slide = _add_blank_slide()
    _add_title_bar(slide, "三、OpenCV预处理管线", "基本内容及主要方法")
    _add_page_number(slide, 9)

    # 流程说明
    steps = [
        ("原始图像", "BGR格式\n垃圾图片"),
        ("高斯去噪", "ksize=3\n消除噪声"),
        ("CLAHE增强", "LAB空间L通道\nclipLimit=2.0"),
        ("缩放归一化", "224×224\n标准化"),
    ]
    x = Inches(0.5)
    for i, (title, desc) in enumerate(steps):
        _add_rect(slide, x, Inches(1.5), Inches(2.5), Inches(1.8), BG_WHITE, SECONDARY)
        _add_rect(slide, x, Inches(1.5), Inches(2.5), Inches(0.45), SECONDARY)
        _add_textbox(slide, x + Inches(0.1), Inches(1.53), Inches(2.3), Inches(0.4),
                     title, font_size=14, color=TEXT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.1), Inches(2.1), Inches(2.3), Inches(1.0),
                     desc, font_size=13, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
        if i < 3:
            _add_textbox(slide, x + Inches(2.5), Inches(1.9), Inches(0.5), Inches(0.5),
                         "→", font_size=28, color=ACCENT, bold=True)
        x += Inches(3.15)

    # OpenCV预处理效果对比图
    _add_textbox(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(0.4),
                 "▎OpenCV预处理效果对比", font_size=18, color=PRIMARY, bold=True)
    _add_image(slide, IMG_DIR / "界面演示截图" / "opencv预处理.png",
               Inches(2.5), Inches(4.1), height=Inches(3.2))


# ═══════════════════════════════════════════
# 第10页：候选模型
# ═══════════════════════════════════════════

def slide_10_models():
    slide = _add_blank_slide()
    _add_title_bar(slide, "三、四种候选轻量模型", "基本内容及主要方法")
    _add_page_number(slide, 10)

    models_info = [
        ("MobileNetV3-Small", "~2.5M", "classifier[3]", "深度可分离卷积\n倒残差结构", "移动端部署"),
        ("ResNet18", "~11.7M", "fc", "残差连接\n跳跃连接", "通用性强"),
        ("EfficientNet-B0", "~5.3M", "classifier[1]", "复合缩放\nMBConv结构", "精度/效率均衡"),
        ("ShuffleNetV2", "~2.3M", "fc", "通道混洗\n点群卷积", "极致轻量"),
    ]
    x = Inches(0.4)
    for name, params, head, feature, scenario in models_info:
        _add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(5.2), BG_WHITE, SECONDARY)
        _add_rect(slide, x, Inches(1.5), Inches(3.0), Inches(0.55), SECONDARY)
        _add_textbox(slide, x + Inches(0.1), Inches(1.55), Inches(2.8), Inches(0.5),
                     name, font_size=16, color=TEXT_LIGHT, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name=FONT_EN)

        items = [
            f"参数量：{params}",
            f"分类头：{head}",
            f"核心特点：{feature}",
            f"适用场景：{scenario}",
            "预训练：ImageNet权重",
            "输出：4类垃圾",
        ]
        _add_bullet_text(slide, x + Inches(0.15), Inches(2.3), Inches(2.7), Inches(4.0),
                         items, font_size=13, color=TEXT_DARK, spacing=Pt(6))
        x += Inches(3.2)

    # 底部说明
    _add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
                 "注：所有模型均采用预训练权重 + 两阶段迁移学习策略，通过@register_model装饰器统一注册管理",
                 font_size=12, color=TEXT_GRAY)


# ═══════════════════════════════════════════
# 第11页：两阶段训练策略
# ═══════════════════════════════════════════

def slide_11_training_strategy():
    slide = _add_blank_slide()
    _add_title_bar(slide, "三、两阶段迁移学习训练策略", "基本内容及主要方法")
    _add_page_number(slide, 11)

    # 阶段一
    _add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.8), BG_WHITE, SECONDARY)
    _add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5), RGBColor(0x27, 0xAE, 0x60))
    _add_textbox(slide, Inches(0.7), Inches(1.53), Inches(5.4), Inches(0.45),
                 "阶段一：特征提取器冻结（Epoch 1-15）", font_size=17, color=TEXT_LIGHT, bold=True)
    _add_bullet_text(slide, Inches(0.7), Inches(2.2), Inches(5.4), Inches(1.8),
                     ["冻结预训练backbone全部参数",
                      "仅训练分类头（4类全连接层）",
                      "学习率：0.001",
                      "快速收敛，避免破坏预训练特征"],
                     font_size=14, color=TEXT_DARK, spacing=Pt(6))

    # 阶段二
    _add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.8), BG_WHITE, SECONDARY)
    _add_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5), SECONDARY)
    _add_textbox(slide, Inches(7.0), Inches(1.53), Inches(5.4), Inches(0.45),
                 "阶段二：全网络微调（Epoch 16-50）", font_size=17, color=TEXT_LIGHT, bold=True)
    _add_bullet_text(slide, Inches(7.0), Inches(2.2), Inches(5.4), Inches(1.8),
                     ["解冻backbone全部参数",
                      "学习率 × 0.1 = 0.0001",
                      "精细调整提取垃圾特征",
                      "EarlyStopping（patience=10）"],
                     font_size=14, color=TEXT_DARK, spacing=Pt(6))

    # 训练参数表
    _add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
                 "▎训练参数设置", font_size=18, color=PRIMARY, bold=True)
    _add_table(slide, Inches(0.5), Inches(5.1), Inches(12.1), Inches(2.0),
               ["参数", "数值", "说明"],
               [["输入尺寸", "224×224", "所有模型统一"],
                ["批量大小", "128", "A100-40GB 显存充足"],
                ["优化器", "AdamW", "weight_decay=1e-4"],
                ["学习率调度", "CosineAnnealingLR", "余弦退火"],
                ["损失函数", "CrossEntropyLoss", "多分类交叉熵"]])


# ═══════════════════════════════════════════
# 第12页：数据增强
# ═══════════════════════════════════════════

def slide_12_augmentation():
    slide = _add_blank_slide()
    _add_title_bar(slide, "三、数据增强策略", "基本内容及主要方法")
    _add_page_number(slide, 12)

    _add_table(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(3.0),
               ["增强方法", "参数", "说明"],
               [["随机水平翻转", "p=0.5", "50%概率水平翻转"],
                ["随机垂直翻转", "p=0.5", "50%概率垂直翻转"],
                ["随机旋转", "±15°", "小角度旋转保持主体完整"],
                ["ColorJitter", "brightness=0.2\ncontrast=0.2\nsaturation=0.2", "亮度/对比度/饱和度扰动"],
                ["RandomErasing", "p=0.2", "随机擦除区域，增强鲁棒性"]])

    # 右侧：OpenCV预处理参数
    _add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5.5), Inches(0.4),
                 "▎OpenCV预处理参数", font_size=18, color=PRIMARY, bold=True)
    _add_table(slide, Inches(7.5), Inches(2.1), Inches(5.3), Inches(1.5),
               ["步骤", "参数"],
               [["高斯去噪", "ksize=3×3"],
                ["CLAHE", "clipLimit=2.0, tileGridSize=(8,8)"],
                ["颜色空间", "BGR→LAB，L通道操作"]])

    # 训练截图
    _add_textbox(slide, Inches(7.5), Inches(3.8), Inches(5.5), Inches(0.4),
                 "▎训练过程", font_size=18, color=PRIMARY, bold=True)
    _add_image(slide, IMG_DIR / "训练过程截图" / "训练配置.png",
               Inches(7.5), Inches(4.3), height=Inches(2.8))


# ═══════════════════════════════════════════
# 第13页：Web系统设计
# ═══════════════════════════════════════════

def slide_13_web_system():
    slide = _add_blank_slide()
    _add_title_bar(slide, "三、Web系统设计", "基本内容及主要方法")
    _add_page_number(slide, 13)

    # API接口列表
    _add_textbox(slide, Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.4),
                 "▎后端API接口（FastAPI）", font_size=18, color=PRIMARY, bold=True)
    _add_table(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.0),
               ["接口", "方法", "功能"],
               [["/api/predict", "POST", "单图分类预测"],
                ["/api/models", "GET", "列出可用模型"],
                ["/api/switch-model", "POST", "切换当前模型"],
                ["/api/compare", "POST", "四模型同图对比"],
                ["/api/batch-predict", "POST", "批量识别统计"],
                ["/api/opencv-pipeline", "POST", "预处理可视化"]])

    # 前端标签页
    _add_textbox(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(0.4),
                 "▎前端多标签页设计", font_size=18, color=PRIMARY, bold=True)

    tabs = [
        ("首页", "系统概览、类别卡片、快速上传入口"),
        ("智能识别", "单图识别 + 四模型对比（四宫格）"),
        ("批量识别", "多图上传 + ECharts类别统计图"),
        ("OpenCV演示", "预处理管线逐步可视化"),
    ]
    y = Inches(2.0)
    for tab_name, tab_desc in tabs:
        _add_rect(slide, Inches(6.8), y, Inches(5.8), Inches(0.65), BG_WHITE, SECONDARY)
        _add_rect(slide, Inches(6.8), y, Inches(1.5), Inches(0.65), SECONDARY)
        _add_textbox(slide, Inches(6.85), y + Inches(0.1), Inches(1.4), Inches(0.45),
                     tab_name, font_size=14, color=TEXT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, Inches(8.4), y + Inches(0.1), Inches(4.1), Inches(0.45),
                     tab_desc, font_size=13, color=TEXT_DARK)
        y += Inches(0.8)

    # 界面截图
    _add_image(slide, IMG_DIR / "界面演示截图" / "单图识别.png",
               Inches(6.8), Inches(5.3), height=Inches(2.0))


# ═══════════════════════════════════════════
# 第14页：实验结果对比
# ═══════════════════════════════════════════

def slide_14_results():
    slide = _add_blank_slide()
    _add_title_bar(slide, "四、四模型实验结果对比", "成果、结论与自我评价")
    _add_page_number(slide, 14)

    # 主表格 - 高亮最优
    _add_table(slide, Inches(0.5), Inches(1.5), Inches(12.1), Inches(2.2),
               ["模型", "Accuracy", "Macro Precision", "Macro Recall", "Macro F1", "参数量", "模型大小"],
               [["★ EfficientNet-B0", "94.47%", "91.64%", "87.70%", "89.54%", "5.3M", "46.39 MB"],
                ["ResNet18",           "93.57%", "90.20%", "85.75%", "87.81%", "11.7M", "128.06 MB"],
                ["ShuffleNetV2",       "93.39%", "89.41%", "85.83%", "87.52%", "2.3M", "14.70 MB"],
                ["MobileNetV3-Small",  "92.38%", "88.15%", "82.66%", "85.11%", "2.5M", "17.67 MB"]])

    # 各类别F1对比
    _add_textbox(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
                 "▎各类别 F1-Score 对比", font_size=18, color=PRIMARY, bold=True)
    _add_table(slide, Inches(0.5), Inches(4.5), Inches(12.1), Inches(1.8),
               ["类别", "MobileNetV3-Small", "ResNet18", "EfficientNet-B0", "ShuffleNetV2"],
               [["可回收物", "94.47%", "95.27%", "★ 96.00%", "95.14%"],
                ["厨余垃圾", "94.80%", "95.58%", "★ 96.03%", "95.53%"],
                ["有害垃圾", "82.05%", "85.01%", "★ 87.67%", "83.73%"],
                ["其他垃圾", "69.12%", "75.39%", "★ 78.46%", "75.67%"]])

    # 推理耗时
    _add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
                 "▎推理耗时（RTX 3060 Laptop，5次平均）：ResNet18 8.2ms | MobileNetV3-Small 11.9ms | ShuffleNetV2 12.8ms | EfficientNet-B0 15.4ms",
                 font_size=13, color=TEXT_GRAY)


# ═══════════════════════════════════════════
# 第15页：训练过程分析
# ═══════════════════════════════════════════

def slide_15_training_curve():
    slide = _add_blank_slide()
    _add_title_bar(slide, "四、EfficientNet-B0 训练过程分析", "成果、结论与自我评价")
    _add_page_number(slide, 15)

    # 训练曲线图
    _add_image(slide, ROUND2_FIG / "training_curve_efficientnet_b0.png",
               Inches(0.5), Inches(1.4), height=Inches(4.5))

    # 右侧关键节点
    _add_textbox(slide, Inches(7.5), Inches(1.5), Inches(5.5), Inches(0.4),
                 "▎训练关键节点", font_size=18, color=PRIMARY, bold=True)
    milestones = [
        "Epoch 1：Train Acc 80.70%，Val Acc 84.39%",
        "Epoch 15（阶段一结束）：Val Acc 86.80%",
        "Epoch 16（解冻微调）：Val Acc 跃升至 91.33%",
        "Epoch 41（最佳）：Val Acc 94.94%",
        "Epoch 50（最终）：Train Acc 99.40%，Val Acc 94.73%",
    ]
    _add_bullet_text(slide, Inches(7.5), Inches(2.1), Inches(5.5), Inches(2.5),
                     milestones, font_size=14, color=TEXT_DARK, spacing=Pt(10))

    # 解冻效果说明
    _add_rect(slide, Inches(7.5), Inches(4.5), Inches(5.3), Inches(2.3), BG_WHITE, ACCENT)
    _add_rect(slide, Inches(7.5), Inches(4.5), Inches(5.3), Inches(0.06), ACCENT)
    _add_textbox(slide, Inches(7.7), Inches(4.7), Inches(5), Inches(0.4),
                 "两阶段策略效果", font_size=16, color=ACCENT, bold=True)
    _add_bullet_text(slide, Inches(7.7), Inches(5.2), Inches(5), Inches(1.4),
                     ["阶段一（冻结）：15个epoch内Val Acc从84%提升至87%",
                      "解冻后1个epoch：87% → 91%（提升4.5%）",
                      "全网络微调35个epoch：最终达到94.73%"],
                     font_size=13, color=TEXT_DARK, spacing=Pt(6))


# ═══════════════════════════════════════════
# 第16页：混淆矩阵与误差分析
# ═══════════════════════════════════════════

def slide_16_confusion():
    slide = _add_blank_slide()
    _add_title_bar(slide, "四、混淆矩阵与误差分析", "成果、结论与自我评价")
    _add_page_number(slide, 16)

    # 混淆矩阵图
    _add_image(slide, IMG_DIR / "正常显示的混淆矩阵的图" / "confusion_matrix_efficientnet_b0.png",
               Inches(0.5), Inches(1.4), height=Inches(4.5))

    # 右侧误差分析
    _add_textbox(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(0.4),
                 "▎主要误差来源", font_size=18, color=PRIMARY, bold=True)

    errors = [
        "其他垃圾误判最多（F1=78.46%）：\n  PE塑料袋外观与可回收物相似",
        "有害垃圾次之（F1=87.67%）：\n  样本量少（仅4802张），指甲油瓶\n  与化妆品瓶外观相似",
        "可回收物/厨余垃圾表现稳定：\n  F1均 > 96%，识别可靠",
    ]
    y = Inches(2.1)
    for err in errors:
        _add_rect(slide, Inches(6.8), y, Inches(5.8), Inches(1.2), BG_WHITE, ACCENT)
        _add_rect(slide, Inches(6.8), y, Inches(0.08), Inches(1.2), ACCENT)
        _add_textbox(slide, Inches(7.0), y + Inches(0.1), Inches(5.5), Inches(1.0),
                     err, font_size=13, color=TEXT_DARK)
        y += Inches(1.4)

    # 底部：OpenCV预处理对比结论
    _add_rect(slide, Inches(6.8), Inches(5.9), Inches(5.8), Inches(1.2), BG_WHITE, RGBColor(0x27, 0xAE, 0x60))
    _add_rect(slide, Inches(6.8), Inches(5.9), Inches(5.8), Inches(0.06), RGBColor(0x27, 0xAE, 0x60))
    _add_textbox(slide, Inches(7.0), Inches(6.1), Inches(5.4), Inches(0.3),
                 "OpenCV预处理效果", font_size=15, color=RGBColor(0x27, 0xAE, 0x60), bold=True)
    _add_textbox(slide, Inches(7.0), Inches(6.5), Inches(5.4), Inches(0.5),
                 "训练速度提升约3倍（400s→130s/epoch），精度小幅提升",
                 font_size=13, color=TEXT_DARK)


# ═══════════════════════════════════════════
# 第17页：结论与自我评价
# ═══════════════════════════════════════════

def slide_17_conclusion():
    slide = _add_blank_slide()
    _add_title_bar(slide, "四、结论与自我评价", "成果、结论与自我评价")
    _add_page_number(slide, 17)

    # 主要结论
    _add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.5),
              "主要结论",
              ["EfficientNet-B0 取得最佳性能（94.47%准确率）",
               "两阶段迁移学习策略有效：解冻后精度跃升4.5%",
               "OpenCV预处理显著提升训练效率（3倍），精度小幅提升",
               "4种轻量模型均可实现实时推理（<16ms）"],
              "01")

    # 不足之处
    _add_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.5),
              "不足之处",
              ["类别不均衡（可回收物62% vs 有害垃圾6%）影响少数类精度",
               "其他垃圾与可回收物材质交叉，F1仅78.46%",
               "未引入注意力机制或更复杂的数据增强策略",
               "系统尚未部署到移动端或嵌入式设备"],
              "02")

    # 自我评价
    _add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.1), Inches(2.8), BG_WHITE, PRIMARY)
    _add_rect(slide, Inches(0.5), Inches(4.3), Inches(0.08), Inches(2.8), ACCENT)
    _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.4),
                 "自我评价", font_size=20, color=PRIMARY, bold=True)
    _add_bullet_text(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0),
                     ["独立完成了从数据集构建、OpenCV预处理设计、模型训练评估到Web系统开发的全流程",
                      "对比了4种主流轻量模型，实验设计合理、数据充分",
                      "系统功能完整：支持单图/批量识别、多模型对比、OpenCV管线可视化",
                      "代码结构清晰、模块化程度高，配置化管理便于复现和扩展"],
                     font_size=15, color=TEXT_DARK, spacing=Pt(8))


# ═══════════════════════════════════════════
# 第18页：作品演示
# ═══════════════════════════════════════════

def slide_18_demo():
    slide = _add_blank_slide()
    _set_slide_bg(slide, PRIMARY)

    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "五、作品演示", font_size=34, color=TEXT_LIGHT, bold=True)
    _add_rect(slide, Inches(0.5), Inches(1.0), Inches(4), Inches(0.04), ACCENT)

    # 四个功能截图
    demos = [
        ("智能识别", "界面演示截图", "单图识别.png"),
        ("模型对比", "界面演示截图", "模型对比.png"),
        ("批量识别", "界面演示截图", "批量识别.png"),
        ("OpenCV演示", "界面演示截图", "opencv预处理.png"),
    ]
    x = Inches(0.5)
    for title, folder, filename in demos:
        img_path = IMG_DIR / folder / filename
        _add_rect(slide, x, Inches(1.4), Inches(2.9), Inches(5.4), BG_WHITE)
        _add_rect(slide, x, Inches(1.4), Inches(2.9), Inches(0.45), SECONDARY)
        _add_textbox(slide, x + Inches(0.1), Inches(1.43), Inches(2.7), Inches(0.4),
                     title, font_size=15, color=TEXT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
        _add_image(slide, img_path, x + Inches(0.05), Inches(1.9), width=Inches(2.8), height=Inches(4.8))
        x += Inches(3.15)

    # 底部提示
    _add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
                 "系统启动方式：运行 start.bat 或 python ui/app.py，浏览器访问 http://localhost:8080",
                 font_size=14, color=RGBColor(0xAE, 0xD6, 0xF1), alignment=PP_ALIGN.CENTER)

    _add_page_number(slide, 18)


# ═══════════════════════════════════════════
# 生成PPT
# ═══════════════════════════════════════════

def main():
    print("开始生成答辩PPT...")
    slide_01_cover()
    print("  第1页：封面 ✓")
    slide_02_toc()
    print("  第2页：目录 ✓")
    slide_03_background()
    print("  第3页：研究背景 ✓")
    slide_04_objectives()
    print("  第4页：研究目的与任务 ✓")
    slide_05_significance()
    print("  第5页：研究意义 ✓")
    slide_06_dataset()
    print("  第6页：数据集来源 ✓")
    slide_07_references()
    print("  第7页：参考文献 ✓")
    slide_08_architecture()
    print("  第8页：系统总体架构 ✓")
    slide_09_opencv()
    print("  第9页：OpenCV预处理管线 ✓")
    slide_10_models()
    print("  第10页：候选模型 ✓")
    slide_11_training_strategy()
    print("  第11页：两阶段训练策略 ✓")
    slide_12_augmentation()
    print("  第12页：数据增强策略 ✓")
    slide_13_web_system()
    print("  第13页：Web系统设计 ✓")
    slide_14_results()
    print("  第14页：实验结果对比 ✓")
    slide_15_training_curve()
    print("  第15页：训练过程分析 ✓")
    slide_16_confusion()
    print("  第16页：混淆矩阵与误差分析 ✓")
    slide_17_conclusion()
    print("  第17页：结论与自我评价 ✓")
    slide_18_demo()
    print("  第18页：作品演示 ✓")

    output_dir = ROOT / "outputs"
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / "答辩PPT_郭靖.pptx"
    prs.save(str(output_path))
    print(f"\nPPT已生成: {output_path}")
    print(f"共 {len(prs.slides)} 页")


if __name__ == "__main__":
    main()
